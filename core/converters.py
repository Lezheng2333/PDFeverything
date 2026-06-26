"""格式转换器 — 各类文件 → PDF 的统一转换接口。

使用 ConverterRegistry 注册表模式：新增格式只需实现 BaseConverter 并注册。"""

import subprocess
import sys
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Callable, List, Optional

from .utils import check_input, read_text_file, temp_pdf_path


# ── AppleScript 模板 ───────────────────────────────────────

_WORD_APPLESCRIPT = '''tell application "Microsoft Word"
    open POSIX file "{input_path}"
    set theDoc to active document
    set pdfPath to POSIX file "{output_path}"
    save as theDoc file name pdfPath file format format PDF
    close theDoc saving no
end tell'''

_PPT_APPLESCRIPT = '''tell application "Microsoft PowerPoint"
    open POSIX file "{input_path}"
    set thePres to active presentation
    save thePres in POSIX file "{output_path}" as save as PDF
    close thePres saving no
end tell'''

_EXCEL_APPLESCRIPT = '''tell application "Microsoft Excel"
    open POSIX file "{input_path}"
    set wb to active workbook
    set ws to active sheet of wb
    tell ws
        set printArea to used range
    end tell
    save wb in POSIX file "{output_path}" as PDF file format
    close wb saving no
end tell'''


def _applescript_convert(app_name: str, script: str, input_path: Path,
                         output_path: Path, timeout: int = 300) -> bool:
    """通过 AppleScript 调用 Office 应用转换文件为 PDF。成功返回 True。"""
    try:
        r = subprocess.run(
            ["osascript", "-e", script.format(
                input_path=str(input_path.resolve()),
                output_path=str(output_path.resolve()),
            )],
            capture_output=True, text=True, timeout=timeout,
        )
        if r.returncode != 0 or not output_path.exists():
            return False
        # 验证输出 PDF 确实有内容（不是空文件）
        try:
            import fitz
            doc = fitz.open(output_path)
            has_pages = len(doc) > 0
            doc.close()
            return has_pages
        except Exception:
            return output_path.stat().st_size > 100  # 至少不是完全的空文件
    except (subprocess.TimeoutExpired, FileNotFoundError, OSError):
        return False


# ── 抽象基类 ────────────────────────────────────────────────


class BaseConverter(ABC):
    """格式转换器抽象基类。"""
    extensions: List[str] = []
    name: str = "Base"

    @abstractmethod
    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        """
        将输入文件转换为 PDF，输出到 output_dir，返回输出的 PDF 路径。
        抛出异常时调用方应捕获并跳过该文件。
        """
        ...


# ── 转换器注册表 ────────────────────────────────────────────


class ConverterRegistry:
    """单例模式：管理所有已注册的格式转换器。"""

    _converters: List[BaseConverter] = []
    _by_ext: dict = {}

    @classmethod
    def register(cls, converter: BaseConverter) -> None:
        cls._converters.append(converter)
        for ext in converter.extensions:
            cls._by_ext[ext.lower()] = converter

    @classmethod
    def get(cls, extension: str) -> Optional[BaseConverter]:
        return cls._by_ext.get(extension.lower().lstrip("."))

    @classmethod
    def supported_extensions(cls) -> List[str]:
        return sorted(cls._by_ext.keys())


# ── 1. 图片转换器 ───────────────────────────────────────────


class ImageConverter(BaseConverter):
    extensions = ["png", "jpg", "jpeg", "gif", "bmp", "tiff", "tif", "webp"]
    name = "图片"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        from PIL import Image

        check_input(input_path)
        img = Image.open(input_path)

        # RGBA → RGB（白色背景）
        if img.mode in ("RGBA", "LA", "P", "PA"):
            background = Image.new("RGB", img.size, (255, 255, 255))
            if img.mode in ("RGBA", "LA", "PA"):
                background.paste(img, mask=img.split()[-1] if img.mode != "P" else None)
            else:
                background.paste(img)
            img = background
        elif img.mode not in ("RGB", "L"):
            img = img.convert("RGB")

        out = output_dir / f"{input_path.stem}.pdf"
        img.save(out, "PDF")
        return out


# ── 2. 文本转换器 ───────────────────────────────────────────


class TextConverter(BaseConverter):
    extensions = ["txt", "md", "log", "py", "json", "xml", "html", "htm",
                  "yaml", "yml", "ini", "cfg", "sh", "bat", "ps1", "csv"]
    name = "文本"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        import fitz

        check_input(input_path)
        content = read_text_file(input_path)
        if not content.strip():
            content = "(空文件)"

        out = output_dir / f"{input_path.stem}.pdf"
        doc = fitz.open()
        page = doc.new_page()
        rect = page.rect

        # 页面参数
        margin = 50
        usable_width = rect.width - 2 * margin
        usable_height = rect.height - 2 * margin
        line_height = 12
        chars_per_line = max(1, int(usable_width / 6.5))  # 等宽字体近似

        y = margin + line_height
        for line in content.split("\n"):
            # 长行截断
            if len(line) > chars_per_line:
                # 按宽度折行
                for chunk_start in range(0, len(line), chars_per_line):
                    chunk = line[chunk_start:chunk_start + chars_per_line]
                    if y + line_height > margin + usable_height:
                        page = doc.new_page()
                        y = margin + line_height
                    page.insert_text(
                        fitz.Point(margin, y), chunk,
                        fontsize=9, fontname="Courier", color=(0, 0, 0),
                    )
                    y += line_height
            else:
                if y + line_height > margin + usable_height:
                    page = doc.new_page()
                    y = margin + line_height
                page.insert_text(
                    fitz.Point(margin, y), line,
                    fontsize=9, fontname="Courier", color=(0, 0, 0),
                )
                y += line_height

        doc.save(out)
        doc.close()
        return out


# ── 3. Word 转换器 ──────────────────────────────────────────


class WordConverter(BaseConverter):
    extensions = ["docx", "doc", "rtf"]
    name = "Word"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        check_input(input_path)
        out = output_dir / f"{input_path.stem}.pdf"

        if progress_callback:
            progress_callback(f"转换 Word: {input_path.name}", 30)

        # 方案 A: AppleScript 调用 Word
        if _applescript_convert("Microsoft Word", _WORD_APPLESCRIPT, input_path, out, timeout=300):
            if progress_callback:
                progress_callback(f"Word 转换完成: {input_path.name}", 100)
            return out

        if progress_callback:
            progress_callback(f"Word 未可用，使用纯 Python 渲染: {input_path.name}", 30)

        # 方案 B: 纯 Python (python-docx + PyMuPDF)
        return self._fallback_convert(input_path, out, progress_callback)

    def _fallback_convert(self, input_path: Path, out: Path,
                          progress_callback=None) -> Path:
        """python-docx 解析 + PyMuPDF 渲染。处理文字、表格和嵌入式图片。"""
        import fitz
        from docx import Document

        doc = Document(input_path)
        pdf = fitz.open()
        page = pdf.new_page()
        rect = page.rect
        margin = 50
        usable_width = rect.width - 2 * margin
        y = margin
        line_height = 15

        def new_page():
            nonlocal y
            pdf.new_page()
            y = margin

        def write_line(text, fontsize=11, bold=False, indent=0):
            nonlocal y
            x = margin + indent
            if y + line_height + 2 > rect.height - margin:
                new_page()
                x = margin + indent
            font = "Helvetica-Bold" if bold else "Helvetica"
            current_page = pdf[-1]
            current_page.insert_text(fitz.Point(x, y), text, fontsize=fontsize, fontname=font)
            y += line_height + 2

        def insert_image_blob(image_blob, max_width=None, max_height=400):
            """将图片数据插入 PDF 当前页面，返回图片高度。"""
            nonlocal y
            if max_width is None:
                max_width = usable_width
            try:
                # 用 PyMuPDF 临时打开图片获取尺寸
                img_doc = fitz.open(stream=image_blob, filetype=None)
                if len(img_doc) == 0:
                    img_doc.close()
                    return 0
                pix = img_doc[0].get_pixmap()
                img_w, img_h = pix.width, pix.height
                img_doc.close()

                # 缩放到适合页面
                if img_w > max_width:
                    scale = max_width / img_w
                    img_w = int(img_w * scale)
                    img_h = int(img_h * scale)
                if img_h > max_height:
                    scale = max_height / img_h
                    img_w = int(img_w * scale)
                    img_h = int(img_h * scale)

                if y + img_h + 5 > rect.height - margin:
                    new_page()

                current_page = pdf[-1]
                img_rect = fitz.Rect(margin, y, margin + img_w, y + img_h)
                # insert_image 直接嵌入光栅图片到 PDF 页面
                current_page.insert_image(img_rect, stream=image_blob)
                y += img_h + 10
                return img_h
            except Exception as e:
                print(f"  [Warning] 图片插入失败: {e}")
                return 0

        # ── 提取段落中的图片（用 python-docx 标准 API） ──
        # namespace 常量
        NS_WP = 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
        NS_A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        NS_R = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'

        for para in doc.paragraphs:
            # 检查段落中是否包含图片
            inline_elems = para._element.findall(f'.//{{{NS_WP}}}inline')
            for inline_elem in inline_elems:
                blips = inline_elem.findall(f'.//{{{NS_A}}}blip')
                for blip in blips:
                    embed = blip.get(f'{{{NS_R}}}embed')
                    if embed:
                        img_part = doc.part.related_parts.get(embed)
                        if img_part and hasattr(img_part, 'blob'):
                            try:
                                insert_image_blob(img_part.blob)
                            except Exception:
                                pass

            # 再处理段落文字
            text = para.text.strip()
            if not text:
                y += line_height // 2
                continue

            # 标题
            if para.style.name.startswith("Heading"):
                fs = {1: 20, 2: 16, 3: 14, 4: 12}.get(
                    int(para.style.name.split()[-1]) if para.style.name.split()[-1].isdigit() else 1, 14)
                write_line(text, fontsize=fs, bold=True)
                y += 4
                continue

            # 缩进
            indent = para.paragraph_format.left_indent
            if indent and hasattr(indent, 'pt'):
                indent = min(int(indent.pt), 200)
            else:
                indent = 0

            is_list = para.style.name.startswith("List") or text.startswith(("•", "-", "*", "·"))

            runs = para.runs
            if runs:
                for run in runs:
                    if run.text.strip():
                        write_line(run.text.strip(), fontsize=11 if not is_list else 10,
                                   bold=run.bold or False, indent=indent)
            else:
                write_line(text, fontsize=11 if not is_list else 10, indent=indent)

        # ── 表格 ──
        for table in doc.tables:
            y += line_height
            write_line("── 表格 ──", fontsize=9, bold=True)
            for row in table.rows:
                cells = [cell.text.replace("\n", " ") for cell in row.cells]
                row_text = " | ".join(cells)
                max_chars = int(usable_width / 6.5)
                if len(row_text) > max_chars:
                    row_text = row_text[:max_chars - 3] + "..."
                write_line(row_text, fontsize=8, bold=False)
            y += 4

        pdf.save(out)
        pdf.close()
        return out


# ── 4. PowerPoint 转换器 ─────────────────────────────────────


class PowerPointConverter(BaseConverter):
    extensions = ["pptx", "ppt"]
    name = "PowerPoint"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        import fitz

        check_input(input_path)
        out = output_dir / f"{input_path.stem}.pdf"

        if progress_callback:
            progress_callback(f"转换 PPT: {input_path.name}", 30)

        # 方案 A: AppleScript
        if _applescript_convert("Microsoft PowerPoint", _PPT_APPLESCRIPT, input_path, out, timeout=300):
            if progress_callback:
                progress_callback(f"PPT 转换完成: {input_path.name}", 100)
            return out

        if progress_callback:
            progress_callback(f"PPT 未可用，使用纯 Python 渲染: {input_path.name}", 30)

        # 方案 B: python-pptx + PyMuPDF
        return self._fallback_convert(input_path, out, progress_callback)

    def _fallback_convert(self, input_path: Path, out: Path,
                          progress_callback=None) -> Path:
        import fitz
        from pptx import Presentation

        prs = Presentation(input_path)
        pdf = fitz.open()

        for slide_idx, slide in enumerate(prs.slides):
            page = pdf.new_page()
            rect = page.rect
            y = 50
            margin = 40

            # 标题
            title = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t and not title:
                            title = t
                            break
                if title:
                    break

            page.insert_text(
                fitz.Point(margin, y),
                f"Slide {slide_idx + 1}: {title}" if title else f"Slide {slide_idx + 1}",
                fontsize=16, fontname="Helvetica-Bold",
            )
            y += 30

            # 提取所有文字
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text or text == title:
                            continue
                        # 折行
                        chars_per_line = max(1, int((rect.width - 2*margin) / 7))
                        for ln in text.split("\n"):
                            if y + 15 > rect.height - margin:
                                page = pdf.new_page()
                                y = 50
                            for chunk_start in range(0, len(ln), chars_per_line):
                                chunk = ln[chunk_start:chunk_start + chars_per_line]
                                page.insert_text(
                                    fitz.Point(margin, y), chunk,
                                    fontsize=10, fontname="Helvetica",
                                )
                                y += 14

            if progress_callback and slide_idx % 5 == 0:
                progress_callback(f"渲染 PPT 第 {slide_idx+1} 页", 30 + int((slide_idx+1)/len(prs.slides)*70))

        pdf.save(out)
        pdf.close()
        return out


# ── 5. Excel 转换器 ─────────────────────────────────────────


class ExcelConverter(BaseConverter):
    extensions = ["xlsx", "xls"]
    name = "Excel"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        check_input(input_path)
        out = output_dir / f"{input_path.stem}.pdf"

        if progress_callback:
            progress_callback(f"转换 Excel: {input_path.name}", 30)

        # 方案 A: AppleScript
        if _applescript_convert("Microsoft Excel", _EXCEL_APPLESCRIPT, input_path, out, timeout=300):
            if progress_callback:
                progress_callback(f"Excel 转换完成: {input_path.name}", 100)
            return out

        if progress_callback:
            progress_callback(f"Excel 未可用，使用纯 Python 渲染: {input_path.name}", 30)

        # 方案 B: openpyxl + PyMuPDF
        return self._fallback_convert(input_path, out, progress_callback)

    def _fallback_convert(self, input_path: Path, out: Path,
                          progress_callback=None) -> Path:
        import fitz
        from openpyxl import load_workbook

        wb = load_workbook(input_path, data_only=True)
        pdf = fitz.open()

        margin = 30
        cell_w = 80   # 每列宽度（点数）
        cell_h = 16   # 行高

        for ws in wb.worksheets:
            page = pdf.new_page()
            rect = page.rect
            usable_width = rect.width - 2 * margin
            max_cols = max(1, int(usable_width / cell_w))

            page.insert_text(
                fitz.Point(margin, 25),
                f"工作表: {ws.title}",
                fontsize=12, fontname="Helvetica-Bold",
            )
            y = 45

            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if row_count > 0 and row_count % 100 == 0:
                    pass  # 后台通知（可选）
                # 截断列数
                row_values = [str(c) if c is not None else "" for c in row][:max_cols]
                if not any(row_values):
                    row_count += 1
                    continue

                if y + cell_h > rect.height - margin:
                    page = pdf.new_page()
                    y = margin

                for ci, val in enumerate(row_values):
                    x = margin + ci * cell_w
                    # 截断过长单元格
                    max_chars = max(1, int(cell_w / 6))
                    display = val[:max_chars - 1] + "…" if len(val) > max_chars else val
                    page.insert_text(
                        fitz.Point(x, y), display,
                        fontsize=8, fontname="Helvetica",
                    )
                y += cell_h
                row_count += 1

        pdf.save(out)
        pdf.close()
        return out


# ── 6. PDF 直通转换器 ───────────────────────────────────────


class PdfPassThroughConverter(BaseConverter):
    """对 PDF 文件直接复制，用于混合合并流水线中的 PDF 项。"""
    extensions = ["pdf"]
    name = "PDF（直通）"

    def convert(self, input_path: Path, output_dir: Path,
                progress_callback: Optional[Callable[[str, int], None]] = None) -> Path:
        import shutil

        check_input(input_path)
        out = output_dir / input_path.name
        shutil.copy2(input_path, out)
        return out


# ── 注册所有转换器 ──────────────────────────────────────────

ConverterRegistry.register(ImageConverter())
ConverterRegistry.register(TextConverter())
ConverterRegistry.register(WordConverter())
ConverterRegistry.register(PowerPointConverter())
ConverterRegistry.register(ExcelConverter())
ConverterRegistry.register(PdfPassThroughConverter())
